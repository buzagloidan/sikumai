import os
import re
import json
import random
import string
import logging
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional

# Import generative AI - using compatible import style
import google.generativeai as genai

# Using more focused libraries for different file types
from unstructured.partition.text import partition_text
# Avoiding unstructured.partition.pdf due to OCR module dependencies
import PyPDF2
import docx
import pptx  # For PowerPoint presentations
from dotenv import load_dotenv

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Set up Gemini AI using compatible configuration method
api_key = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=api_key)

class QuestionGenerator:
    """Generate quiz questions from text content using Gemini 2.0 Flash."""
    
    def __init__(self):
        self.gemini_model = "gemini-2.0-flash"  # Using Gemini 2.0 Flash
    
    def extract_text(self, file_content: bytes, mime_type: str) -> str:
        """Extract text from file using appropriate libraries based on file type."""
        try:
            # Create a temporary file to work with
            suffix = self._get_file_suffix(mime_type)
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            text_content = ""
            
            # Use appropriate extraction method based on file type
            if mime_type == 'application/pdf':
                # Use PyPDF2 for PDF files
                with open(temp_file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        text_content += page.extract_text() + "\n"
                        
            elif mime_type in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                # Use python-docx for DOCX files
                doc = docx.Document(temp_file_path)
                for para in doc.paragraphs:
                    text_content += para.text + "\n"
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                # Use python-pptx for PPTX files
                presentation = pptx.Presentation(temp_file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                    # Add a separator between slides
                    text_content += "\n---\n"
                    
            elif mime_type == 'text/plain':
                # Simple text file reading
                with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text_content = f.read()
            else:
                # Fallback for unsupported mime types - try as text
                try:
                    with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text_content = f.read()
                except:
                    logger.error(f"Unsupported file type: {mime_type}")
                    raise ValueError(f"Unsupported file type: {mime_type}")
            
            # Clean up the temporary file
            os.unlink(temp_file_path)
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            raise ValueError(f"Failed to extract text: {str(e)}")
    
    def _get_file_suffix(self, mime_type: str) -> str:
        """Get file suffix based on MIME type."""
        mime_to_suffix = {
            'application/pdf': '.pdf',
            'application/msword': '.doc',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'text/plain': '.txt',
            'image/jpeg': '.jpg',
            'image/png': '.png'
        }
        return mime_to_suffix.get(mime_type, '.tmp')
    
    def clean_text(self, text: str) -> str:
        """Clean and preprocess text."""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove non-printable characters
        text = "".join(c if c.isprintable() or c in ['\n', '\t'] else ' ' for c in text)
        return text.strip()
    
    def generate_questions(self, file_content: bytes, mime_type: str, num_questions: int = 20) -> List[Dict]:
        """
        Generate quiz questions from file content using Gemini 2.0 Flash's large context window.
        
        Args:
            file_content: Binary content of the file
            mime_type: MIME type of the file
            num_questions: Number of questions to generate (always 20 as per requirements)
            
        Returns:
            List of question dictionaries - always 20 questions
        """
        try:
            # Extract text from file using the appropriate method
            text = self.extract_text(file_content, mime_type)
            clean_text = self.clean_text(text)
            
            # If text is too short, return an error
            if len(clean_text) < 100:
                logger.error("Extracted text is too short")
                raise ValueError("The document contains too little text to generate questions")
            
            logger.info(f"Generating {num_questions} questions with Gemini 2.0 Flash")
            
            # Configure the generation parameters for Gemini 2.0 Flash
            generation_config = {
                "temperature": round(random.uniform(0.9, 1.0), 2),  # Randomize temperature for diversity
                "top_p": 1,
                "top_k": 32,
                "max_output_tokens": 8192,
            }
            
            logger.info(f"Using temperature: {generation_config['temperature']}")
            
            # Set safety settings to lowest level - BLOCK_NONE for all categories
            safety_settings = [
                {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
            ]
            
            # Utilize Gemini's large context window (up to 1M tokens)
            # We'll use 200K characters which is a safe limit while still being much larger than before
            max_content_length = 200000
            
            if len(clean_text) > max_content_length:
                logger.warning(f"Content length ({len(clean_text)}) exceeds maximum ({max_content_length}), truncating")
                content_for_prompt = clean_text[:max_content_length]
            else:
                content_for_prompt = clean_text
            
            # Create the prompt to generate all questions at once
            prompt = f"""
            Create EXACTLY 20 multiple choice questions in Hebrew that assess mastery 
            of the concepts from the background content. Generate questions that could be answered by someone 
            who truly understands the material, without needing to reference specific text.

            THE NUMBER OF QUESTIONS MUST BE EXACTLY 20. THIS IS CRITICAL.

            CRITICAL RULES:
            1. NEVER use phrases like 'according to the text', 'based on the passage', or any direct text references
            2. Questions must be in Hebrew
            3. Each question must have exactly 4 options WITHOUT any prefixes or labels
            4. The correct answer must be unambiguously correct and fully supported by the background content
            5. Focus on testing:
               - Deep comprehension of concepts
               - Ability to apply principles
               - Understanding of relationships and implications
               - Critical thinking about the subject matter
            6. Each explanation must clearly justify why the correct answer is the only valid choice
            7. Do not use trailing commas in arrays
            8. The questions should cover different aspects of the document
            9. EXACTLY 20 QUESTIONS - NO MORE, NO LESS
            10. IMPORTANT: All 4 answer options must be of approximately equal length and complexity
            11. All answer options must be plausible to avoid obvious wrong options
            12. Don't make the correct answer more detailed or longer than incorrect options

            Return a valid JSON array where each question has this exact format:
            {{
                "question": "שאלה בעברית?",
                "options": ["אפשרות 1", "אפשרות 2", "אפשרות 3", "אפשרות 4"],
                "correct_option_index": 0,
                "explanation": "הסבר קצר"
            }}

            Background content to derive concepts from:
            {content_for_prompt}
            """
            
            # Make the API call with retries
            max_attempts = 3
            all_questions = []
            attempt = 0
            
            # Create Gemini model instance
            model = genai.GenerativeModel(
                model_name=self.gemini_model,
                generation_config=generation_config,
                safety_settings=safety_settings
            )
            
            # Try to generate all questions in one go
            while attempt < max_attempts and len(all_questions) < num_questions:
                attempt += 1
                logger.warning(f"Attempt {attempt} to generate all questions")
                
                try:
                    # Generate content using compatible API format
                    response = model.generate_content(prompt)
                    
                    if not response or not hasattr(response, 'text'):
                        logger.warning("Empty response from Gemini API")
                        continue
                        
                    response_text = response.text
                    
                    # Clean the response to ensure it's valid JSON
                    response_text = re.sub(r'^```json', '', response_text)
                    response_text = re.sub(r'```$', '', response_text)
                    response_text = re.sub(r'^```', '', response_text)
                    response_text = response_text.strip()
                    
                    # Try to parse JSON
                    try:
                        questions = json.loads(response_text)
                        
                        # Ensure questions is always a list
                        if not isinstance(questions, list):
                            questions = [questions]
                        
                        logger.warning(f"Successfully parsed JSON response with {len(questions)} questions")
                        
                        # Validate and process questions
                        processed_questions = []
                        for q in questions:
                            # Check for required fields with possible field name variations
                            if 'correct_option_index' in q and 'correctAnswer' not in q:
                                q['correctAnswer'] = q['correct_option_index']
                            
                            # Validate the question format
                            valid = True
                            if not all(key in q for key in ['question', 'options']):
                                logger.warning(f"Question missing required fields: {q}")
                                valid = False
                                
                            if 'correctAnswer' not in q and 'correct_option_index' not in q:
                                logger.warning(f"Question missing correct answer index: {q}")
                                valid = False
                                
                            if len(q.get('options', [])) != 4:
                                logger.warning(f"Question does not have exactly 4 options: {q}")
                                valid = False
                                
                            correct_idx = q.get('correctAnswer', q.get('correct_option_index', -1))
                            if not isinstance(correct_idx, int) or correct_idx not in range(4):
                                logger.warning(f"correctAnswer must be an integer between 0-3: {q}")
                                valid = False
                            
                            if valid:
                                # Standardize field names
                                correct_idx = q.get('correctAnswer', q.get('correct_option_index', 0))
                                options = q['options']
                                correct_option = options[correct_idx]
                                
                                # Randomize the position of the correct answer
                                shuffled_options = options.copy()
                                random.shuffle(shuffled_options)
                                new_correct_idx = shuffled_options.index(correct_option)
                                
                                logger.info(f"Randomized options: original correct idx={correct_idx}, new correct idx={new_correct_idx}")
                                
                                processed_question = {
                                    'id': ''.join(random.choices(string.ascii_lowercase + string.digits, k=10)),
                                    'question': q['question'],
                                    'options': shuffled_options,
                                    'correctAnswer': new_correct_idx,
                                    'explanation': q.get('explanation', '')
                                }
                                processed_questions.append(processed_question)
                        
                        all_questions.extend(processed_questions)
                        
                        # If we got sufficient questions, break
                        if len(all_questions) >= num_questions:
                            break
                            
                    except json.JSONDecodeError as e:
                        logger.error(f"Failed to parse JSON response: {e}")
                        logger.error(f"Response text: {response_text}")
                except Exception as e:
                    logger.error(f"Error generating questions: {e}")
            
            # If we still don't have enough questions, generate them one by one
            if len(all_questions) < num_questions:
                logger.warning(f"Only generated {len(all_questions)} questions in batch mode, generating remaining individually")
                remaining = num_questions - len(all_questions)
                
                # Generate individual questions using smaller chunks of the content
                chunk_size = len(content_for_prompt) // remaining
                for i in range(remaining):
                    try:
                        # Get a chunk of the content
                        start_idx = (i * chunk_size) % max(1, len(content_for_prompt) - chunk_size)
                        chunk = content_for_prompt[start_idx:start_idx + chunk_size]
                        
                        # Randomize temperature for this individual question
                        individual_temp = round(random.uniform(0.9, 1.0), 2)
                        model.generation_config["temperature"] = individual_temp
                        logger.info(f"Using temperature {individual_temp} for individual question #{i+1}")
                        
                        # Create prompt for a single question
                        single_prompt = f"""
                        Create exactly 1 multiple choice question in Hebrew that assesses mastery 
                        of the concepts from the background content. The question should test understanding 
                        of the material without directly referencing the text.

                        CRITICAL RULES:
                        1. Question must be in Hebrew
                        2. Must have exactly 4 options WITHOUT any prefixes
                        3. The correct answer must be unambiguously correct
                        4. IMPORTANT: All 4 answer options must be of approximately equal length and complexity
                        5. The correct answer should NOT be more detailed or longer than incorrect options
                        6. All answer options must be plausible and look legitimate

                        Return ONLY ONE question in this JSON format:
                        {{
                            "question": "שאלה בעברית?",
                            "options": ["אפשרות 1", "אפשרות 2", "אפשרות 3", "אפשרות 4"],
                            "correct_option_index": 0,
                            "explanation": "הסבר קצר"
                        }}

                        Background content:
                        {chunk}
                        """
                        
                        logger.warning(f"Generating individual question #{i+1}")
                        
                        # Generate individual question using compatible API format
                        # We can reuse the model instance from above
                        response = model.generate_content(single_prompt)
                        
                        if response and hasattr(response, 'text'):
                            # Clean the response
                            response_text = response.text
                            response_text = re.sub(r'^```json', '', response_text)
                            response_text = re.sub(r'```$', '', response_text)
                            response_text = re.sub(r'^```', '', response_text)
                            response_text = response_text.strip()
                            
                            try:
                                question_data = json.loads(response_text)
                            except:
                                logger.error(f"Failed to parse individual question response: {response_text}")
                                continue
                            
                            if question_data and 'question' in question_data and 'options' in question_data:
                                # Get correct index
                                correct_idx = question_data.get('correctAnswer', 
                                             question_data.get('correct_option_index', 0))
                                
                                options = question_data['options'][:4]  # Ensure exactly 4 options
                                
                                # Avoid index errors
                                if not 0 <= correct_idx < len(options):
                                    correct_idx = 0
                                    
                                correct_option = options[correct_idx]
                                
                                # Randomize the position of the correct answer
                                shuffled_options = options.copy()
                                random.shuffle(shuffled_options)
                                new_correct_idx = shuffled_options.index(correct_option)
                                
                                logger.info(f"Individual question: original correct idx={correct_idx}, new correct idx={new_correct_idx}")
                                
                                # Create standardized question
                                processed_question = {
                                    'id': ''.join(random.choices(string.ascii_lowercase + string.digits, k=10)),
                                    'question': question_data['question'],
                                    'options': shuffled_options,
                                    'correctAnswer': new_correct_idx,
                                    'explanation': question_data.get('explanation', '')
                                }
                                
                                all_questions.append(processed_question)
                                logger.warning(f"Successfully generated individual question #{i+1}")
                    except Exception as e:
                        logger.error(f"Error generating individual question #{i+1}: {e}")
            
            # Final validation - ensure we have exactly the right number of questions
            if len(all_questions) > num_questions:
                # Trim to the exact number needed
                all_questions = all_questions[:num_questions]
            
            logger.warning(f"Final question count: {len(all_questions)}")
            
            return all_questions
            
        except Exception as e:
            logger.error(f"Critical error in generate_questions: {str(e)}")
            raise e  # Re-raise the exception to be handled by the caller

# Example usage
if __name__ == "__main__":
    generator = QuestionGenerator()
    # Example would load a file and generate questions
    print("Question generator ready - using specialized parsers and Gemini for PDF, DOCX, PPTX and TXT files") 